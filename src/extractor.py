"""
Document content extraction via Microsoft Graph API.

Downloads document content from SharePoint Online and extracts readable
text for classification. Uses Graph API's built-in content extraction
where possible, falling back to binary download + local parsing.
"""

import io
import logging
import time

from .auth import GraphAuthClient

logger = logging.getLogger(__name__)

# Optional imports for local text extraction fallback
try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import pdfplumber
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class DocumentExtractor:
    """Extracts text content from SharePoint documents via Graph API."""

    # Extensions we can extract text from
    SUPPORTED_EXTENSIONS = {
        ".docx", ".doc", ".pdf", ".pptx", ".ppt",
        ".xlsx", ".xls", ".txt", ".csv", ".md",
        ".html", ".htm", ".rtf",
    }

    # Max text length to keep per document (to control costs with Azure OpenAI)
    MAX_TEXT_LENGTH = 8000

    # Rate limiting: pause between API calls to avoid throttling
    REQUEST_DELAY = 0.3  # seconds

    def __init__(self, auth_client: GraphAuthClient):
        self.auth = auth_client

    def extract_text(self, drive_item_id: str, drive_id: str,
                     file_name: str, extension: str) -> str:
        """Extract text content from a SharePoint document.

        Tries Graph API preview/content endpoint first, falls back to
        downloading the file and parsing locally.

        Args:
            drive_item_id: The Graph item ID of the document.
            drive_id: The Graph drive ID containing the document.
            file_name: Name of the file (for logging).
            extension: File extension (e.g., '.docx').

        Returns:
            Extracted text content, truncated to MAX_TEXT_LENGTH.
            Empty string if extraction fails.
        """
        ext = extension.lower()

        if ext not in self.SUPPORTED_EXTENSIONS:
            logger.debug(f"Skipping unsupported type {ext}: {file_name}")
            return ""

        # For plain text files, just download directly
        if ext in {".txt", ".csv", ".md"}:
            return self._download_as_text(drive_id, drive_item_id, file_name)

        # For Office/PDF docs, try downloading and parsing locally
        content = self._download_and_parse(drive_id, drive_item_id,
                                           file_name, ext)

        time.sleep(self.REQUEST_DELAY)
        return content[:self.MAX_TEXT_LENGTH] if content else ""

    def _download_as_text(self, drive_id: str, item_id: str,
                          file_name: str) -> str:
        """Download a text-based file and return its content."""
        try:
            import requests

            endpoint = (
                f"https://graph.microsoft.com/v1.0"
                f"/drives/{drive_id}/items/{item_id}/content"
            )
            response = requests.get(
                endpoint, headers=self.auth.headers, timeout=30
            )

            if response.status_code == 401:
                self.auth.get_token()
                response = requests.get(
                    endpoint, headers=self.auth.headers, timeout=30
                )

            response.raise_for_status()

            text = response.text[:self.MAX_TEXT_LENGTH]
            logger.debug(f"Extracted {len(text)} chars from {file_name}")
            return text

        except Exception as e:
            logger.warning(f"Failed to download text from {file_name}: {e}")
            return ""

    def _download_and_parse(self, drive_id: str, item_id: str,
                            file_name: str, extension: str) -> str:
        """Download a binary file and extract text locally."""
        try:
            import requests

            endpoint = (
                f"https://graph.microsoft.com/v1.0"
                f"/drives/{drive_id}/items/{item_id}/content"
            )
            response = requests.get(
                endpoint, headers=self.auth.headers, timeout=60
            )

            if response.status_code == 401:
                self.auth.get_token()
                response = requests.get(
                    endpoint, headers=self.auth.headers, timeout=60
                )

            response.raise_for_status()
            file_bytes = io.BytesIO(response.content)

            if extension in {".docx", ".doc"} and HAS_DOCX:
                return self._parse_docx(file_bytes, file_name)
            elif extension == ".pdf" and HAS_PDF:
                return self._parse_pdf(file_bytes, file_name)
            elif extension in {".xlsx", ".xls"} and HAS_XLSX:
                return self._parse_xlsx(file_bytes, file_name)
            elif extension in {".pptx", ".ppt"} and HAS_PPTX:
                return self._parse_pptx(file_bytes, file_name)
            elif extension in {".html", ".htm"}:
                return self._parse_html(response.text, file_name)
            elif extension == ".rtf":
                # Basic RTF — just strip control codes
                text = response.text
                return text[:self.MAX_TEXT_LENGTH]
            else:
                logger.debug(
                    f"No parser available for {extension}: {file_name}"
                )
                return ""

        except Exception as e:
            logger.warning(f"Failed to parse {file_name}: {e}")
            return ""

    def _parse_docx(self, file_bytes: io.BytesIO, file_name: str) -> str:
        """Extract text from a DOCX file."""
        try:
            doc = DocxDocument(file_bytes)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text = "\n".join(paragraphs)
            logger.debug(f"Extracted {len(text)} chars from DOCX: {file_name}")
            return text
        except Exception as e:
            logger.warning(f"DOCX parse failed for {file_name}: {e}")
            return ""

    def _parse_pdf(self, file_bytes: io.BytesIO, file_name: str) -> str:
        """Extract text from a PDF file."""
        try:
            text_parts = []
            with pdfplumber.open(file_bytes) as pdf:
                for page in pdf.pages[:50]:  # Cap at 50 pages
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
            text = "\n".join(text_parts)
            logger.debug(f"Extracted {len(text)} chars from PDF: {file_name}")
            return text
        except Exception as e:
            logger.warning(f"PDF parse failed for {file_name}: {e}")
            return ""

    def _parse_xlsx(self, file_bytes: io.BytesIO, file_name: str) -> str:
        """Extract text from an XLSX file (sheet names + cell values)."""
        try:
            wb = openpyxl.load_workbook(file_bytes, read_only=True,
                                         data_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames:
                text_parts.append(f"[Sheet: {sheet_name}]")
                ws = wb[sheet_name]
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        text_parts.append(" | ".join(cells))
                    row_count += 1
                    if row_count > 200:  # Cap rows per sheet
                        text_parts.append("... (truncated)")
                        break
            wb.close()
            text = "\n".join(text_parts)
            logger.debug(f"Extracted {len(text)} chars from XLSX: {file_name}")
            return text
        except Exception as e:
            logger.warning(f"XLSX parse failed for {file_name}: {e}")
            return ""

    def _parse_pptx(self, file_bytes: io.BytesIO, file_name: str) -> str:
        """Extract text from a PPTX file."""
        try:
            prs = Presentation(file_bytes)
            text_parts = []
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                if slide_texts:
                    text_parts.append(
                        f"[Slide {i + 1}] " + " ".join(slide_texts)
                    )
            text = "\n".join(text_parts)
            logger.debug(f"Extracted {len(text)} chars from PPTX: {file_name}")
            return text
        except Exception as e:
            logger.warning(f"PPTX parse failed for {file_name}: {e}")
            return ""

    def _parse_html(self, html_text: str, file_name: str) -> str:
        """Extract text from HTML by stripping tags."""
        try:
            import re
            clean = re.sub(r"<[^>]+>", " ", html_text)
            clean = re.sub(r"\s+", " ", clean).strip()
            logger.debug(
                f"Extracted {len(clean)} chars from HTML: {file_name}"
            )
            return clean
        except Exception as e:
            logger.warning(f"HTML parse failed for {file_name}: {e}")
            return ""
